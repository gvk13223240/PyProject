import streamlit as st
from pptx import Presentation
from pptx.util import Inches
from PIL import Image
import io

# Layout options
layout_options = {
    "1 image per slide (1x1)": (1, 1),
    "2 images per slide (2x1)": (2, 1),
    "4 images per slide (2x2)": (2, 2),
    "6 images per slide (3x2)": (3, 2)
}

# Session state setup
if "preview_mode" not in st.session_state:
    st.session_state.preview_mode = False
if "uploaded_files" not in st.session_state:
    st.session_state.uploaded_files = None
if "layout" not in st.session_state:
    st.session_state.layout = None

# Title
st.title("📸 Images to PPTX Converter")

# Upload and layout selection
if not st.session_state.preview_mode:
    uploaded_files = st.file_uploader("Upload images", type=['png', 'jpg', 'jpeg'], accept_multiple_files=True)
    layout_name = st.selectbox("Choose layout", list(layout_options.keys()))

    if uploaded_files:
        st.session_state.uploaded_files = uploaded_files
        st.session_state.layout = layout_options[layout_name]

        if st.button("🔍 Preview Slides"):
            st.session_state.preview_mode = True
            st.experimental_rerun()

# Slide preview
if st.session_state.preview_mode:
    def create_slide_mockups(files, layout):
        previews = []
        cols, rows = layout
        per_slide = cols * rows
        canvas_w, canvas_h = 200 * cols, 150 * rows
        box_w, box_h = canvas_w // cols, canvas_h // rows

        for i in range(0, len(files), per_slide):
            canvas = Image.new("RGB", (canvas_w, canvas_h), "white")
            for j, f in enumerate(files[i:i+per_slide]):
                img = Image.open(f)
                img.thumbnail((box_w, box_h))
                x = (j % cols) * box_w
                y = (j // cols) * box_h
                canvas.paste(img, (x, y))
            previews.append(canvas)
        return previews

    st.subheader("📽️ Slide Previews")
    slides = create_slide_mockups(st.session_state.uploaded_files, st.session_state.layout)
    for i, img in enumerate(slides, 1):
        st.image(img, caption=f"Slide {i}", use_column_width=True)

    def generate_pptx(files, layout, filename="Screenshots_Presentation.pptx"):
        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        sw, sh = prs.slide_width, prs.slide_height
        cols, rows = layout
        per_slide = cols * rows
        img_w, img_h = sw // cols, sh // rows

        for i in range(0, len(files), per_slide):
            slide = prs.slides.add_slide(blank_layout)
            for j, f in enumerate(files[i:i+per_slide]):
                img = Image.open(f)
                stream = io.BytesIO()
                img.save(stream, format="PNG")
                stream.seek(0)
                x = (j % cols) * img_w
                y = (j // cols) * img_h
                slide.shapes.add_picture(stream, x, y, width=img_w, height=img_h)
        prs.save(filename)

    if st.button("📥 Download PPTX"):
        generate_pptx(st.session_state.uploaded_files, st.session_state.layout)
        with open("Screenshots_Presentation.pptx", "rb") as f:
            st.download_button("Download Final PPTX", f, file_name="Screenshots_Presentation.pptx", mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")

    if st.button("⬅️ Go Back"):
        st.session_state.preview_mode = False
        st.experimental_rerun()

# Footer
st.markdown("""
    <style>
    .footer {
        position: fixed;
        left: 0;
        bottom: 0;
        width: 100%;
        text-align: center;
        color: gray;
        font-size: 12px;
        padding: 10px 0;
    }
    </style>
    <div class="footer">
        &copy; 2025 Vamshi Krishna G. All rights reserved.
    </div>
""", unsafe_allow_html=True)
