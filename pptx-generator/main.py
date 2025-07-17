import streamlit as st
from pptx import Presentation
from pptx.util import Inches
from PIL import Image
import io

st.set_page_config(page_title="Images to PPTX", layout="centered")

def screenshots_to_pptx_from_memory(image_files, output_path="Screenshots_Presentation.pptx", layout=(1, 1)):
    prs = Presentation()
    blank_layout = prs.slide_layouts[6]
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    cols, rows = layout
    per_slide = cols * rows
    img_index = 0

    for i in range(0, len(image_files), per_slide):
        slide = prs.slides.add_slide(blank_layout)
        for row in range(rows):
            for col in range(cols):
                if img_index >= len(image_files):
                    break
                image_file = image_files[img_index]
                img = Image.open(image_file)
                img_stream = io.BytesIO()
                img.save(img_stream, format="PNG")
                img_stream.seek(0)

                box_width = slide_width / cols
                box_height = slide_height / rows

                left = Inches(col * (box_width / Inches(1)))
                top = Inches(row * (box_height / Inches(1)))
                slide.shapes.add_picture(img_stream, left, top, width=box_width, height=box_height)

                img_index += 1

    prs.save(output_path)


def create_slide_mockups(files, layout):
    previews = []
    cols, rows = layout
    per_slide = cols * rows
    box_w, box_h = 400, 300
    canvas_w, canvas_h = box_w * cols, box_h * rows

    for i in range(0, len(files), per_slide):
        canvas = Image.new("RGB", (canvas_w, canvas_h), "white")
        for j, f in enumerate(files[i:i+per_slide]):
            try:
                img = Image.open(f).convert("RGB")
                img = img.resize((box_w, box_h), resample=Image.Resampling.LANCZOS)
                x = (j % cols) * box_w
                y = (j // cols) * box_h
                canvas.paste(img, (x, y))
            except Exception as e:
                print(f"Error loading image {f.name}: {e}")
        previews.append(canvas)
    return previews


# --- UI ---
st.title("📸 Images to PPTX Converter")

uploaded_files = st.file_uploader("Upload images", accept_multiple_files=True, type=["png", "jpg", "jpeg"])

layout_option = st.selectbox("Choose images per slide layout", options=["1x1", "2x1", "2x2", "3x2"])
layout_map = {"1x1": (1, 1), "2x1": (2, 1), "2x2": (2, 2), "3x2": (3, 2)}
layout = layout_map[layout_option]

if uploaded_files:
    if st.button("🎯 Generate PPTX"):
        screenshots_to_pptx_from_memory(uploaded_files, layout=layout)
        with open("Screenshots_Presentation.pptx", "rb") as f:
            st.download_button(
                label="⬇️ Download PPTX",
                data=f,
                file_name="Screenshots_Presentation.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )
        st.success("PPTX generated successfully!")
        st.balloons()

    st.markdown("---")
    st.subheader("🔍 Preview Slides Before Download")
    previews = create_slide_mockups(uploaded_files, layout)

    for i, preview in enumerate(previews):
        st.image(preview, caption=f"Slide {i+1}", use_container_width=True)

# Footer
st.markdown(
    """
    <style>
    .footer {
        text-align: center;
        color: gray;
        font-size: 12px;
        padding: 10px 0;
        margin-top: 30px;
    }
    </style>
    <div class="footer">
        &copy; 2025 Vamshi Krishna G. All rights reserved.
    </div>
    """,
    unsafe_allow_html=True
)
